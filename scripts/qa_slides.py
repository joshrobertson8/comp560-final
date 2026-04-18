"""
Structural QA for the slide deck (no LibreOffice required).

Flags likely overflow issues by comparing estimated text extent against its
textbox dimensions.

Heuristic: an average word at N pt takes roughly (N * 0.55) pt of horizontal
space per character (for Calibri-style sans-serif). A line of body text at
15 pt is ~18 pt tall including spacing.
"""
from pathlib import Path

from pptx import Presentation
from pptx.util import Emu

ROOT = Path(__file__).resolve().parent.parent
PPTX = ROOT / "joshrobertson_comp560_slides.pptx"


def emu_to_in(emu):
    return emu / 914400.0


def iter_textframes(slide):
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        tf = shape.text_frame
        yield shape, tf


def estimate_lines(text: str, width_in: float, font_pt: float) -> int:
    """Rough estimate of wrapped-line count."""
    if not text:
        return 0
    char_w_pt = font_pt * 0.55
    width_pt = width_in * 72
    chars_per_line = max(1, int(width_pt / char_w_pt))
    lines = 0
    for line in text.split("\n"):
        n = len(line)
        lines += max(1, -(-n // chars_per_line))  # ceil
    return lines


def main():
    prs = Presentation(str(PPTX))
    print(f"Slide size: {emu_to_in(prs.slide_width):.2f} × {emu_to_in(prs.slide_height):.2f} in")
    print(f"{len(prs.slides)} slides\n")

    issues = []
    for i, slide in enumerate(prs.slides, 1):
        print(f"--- Slide {i} ---")
        for shape, tf in iter_textframes(slide):
            w_in = emu_to_in(shape.width)
            h_in = emu_to_in(shape.height)
            x_in = emu_to_in(shape.left)
            y_in = emu_to_in(shape.top)

            # Collect all runs and their font sizes
            full_text = ""
            max_pt = 0
            for p in tf.paragraphs:
                line = "".join(r.text for r in p.runs)
                full_text += line + "\n"
                for r in p.runs:
                    if r.font.size:
                        max_pt = max(max_pt, r.font.size.pt)
            full_text = full_text.rstrip("\n")
            if not full_text.strip():
                continue

            pt = max_pt or 14
            lines = estimate_lines(full_text, w_in, pt)
            line_h_in = (pt * 1.25) / 72.0  # ~1.25 line-height
            estimated_h = lines * line_h_in

            # Right-edge overflow (text extending past slide width)
            if x_in + w_in > emu_to_in(prs.slide_width) + 0.05:
                issues.append(f"Slide {i}: textbox extends past slide width — '{full_text[:40]}...'")

            # Height overflow
            if estimated_h > h_in + 0.15:
                issues.append(
                    f"Slide {i}: possible height overflow — box {w_in:.2f}x{h_in:.2f}, "
                    f"text needs ~{estimated_h:.2f}in ({lines} lines @ {pt:.0f}pt) — '{full_text[:50]}...'"
                )

            # Bottom edge
            if y_in + h_in > emu_to_in(prs.slide_height) + 0.05:
                issues.append(f"Slide {i}: textbox extends past slide bottom — '{full_text[:40]}...'")

            print(f"  [{pt:>3.0f}pt] {w_in:4.2f}x{h_in:4.2f}in  est {lines}L / {estimated_h:.2f}in  "
                  f"({full_text[:60].replace(chr(10), ' ⏎ ')!r})")

    print()
    if issues:
        print("=" * 60)
        print(f"ISSUES FOUND ({len(issues)}):")
        for msg in issues:
            print(f"  - {msg}")
    else:
        print("No overflow risks flagged.")


if __name__ == "__main__":
    main()
