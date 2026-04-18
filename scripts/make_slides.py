"""
Build the final slide deck for the COMP 560 project.

10 slides, 16:9, dark/light sandwich layout, cohesive palette with the charts.
"""
from __future__ import annotations

from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

ROOT = Path(__file__).resolve().parent.parent
CHARTS = ROOT / "charts"
OUT = ROOT / "joshrobertson_comp560_slides.pptx"

# ---- Palette (cohesive with matplotlib charts) ----
PRIMARY = RGBColor(0x2E, 0x86, 0xAB)       # chart primary blue
DARK = RGBColor(0x1A, 0x3A, 0x52)          # deep teal for title slides
DARKER = RGBColor(0x0F, 0x26, 0x38)        # nearly-black accent
ACCENT = RGBColor(0xE6, 0x39, 0x46)        # chart accent red/coral
LIGHT = RGBColor(0xF5, 0xF7, 0xFA)         # page bg for content slides
TEXT_DARK = RGBColor(0x1A, 0x1A, 0x1A)
TEXT_MUTED = RGBColor(0x5A, 0x6C, 0x78)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)

# ---- Setup ----
prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
SW, SH = prs.slide_width, prs.slide_height

BLANK = prs.slide_layouts[6]  # blank layout


# ---------------------------------------------------------------------------
# Helpers
# ---------------------------------------------------------------------------

def fill(shape, color: RGBColor):
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()


def add_rect(slide, x, y, w, h, color: RGBColor):
    r = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    fill(r, color)
    return r


def add_text(slide, x, y, w, h, text, *, size=16, bold=False, color=TEXT_DARK,
             align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, font="Calibri",
             line_spacing=1.15):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.margin_left = Inches(0)
    tf.margin_right = Inches(0)
    tf.margin_top = Inches(0)
    tf.margin_bottom = Inches(0)
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    p = tf.paragraphs[0]
    p.alignment = align
    p.line_spacing = line_spacing
    r = p.add_run()
    r.text = text
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.color.rgb = color
    r.font.name = font
    return tb, p, r


def add_bullets(slide, x, y, w, h, items, *, size=16, color=TEXT_DARK,
                bullet="•", line_spacing=1.35, indent_first=0.25):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.margin_left = Inches(0)
    tf.margin_right = Inches(0)
    tf.margin_top = Inches(0)
    tf.margin_bottom = Inches(0)
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        p.line_spacing = line_spacing
        p.space_after = Pt(6)
        r = p.add_run()
        r.text = f"{bullet}  {item}"
        r.font.size = Pt(size)
        r.font.color.rgb = color
        r.font.name = "Calibri"
    return tb


def add_page_header(slide, title: str, kicker: str | None = None):
    """Light-bg content slide header: kicker + title + accent line.
    Title box is tall enough to accommodate 2-line wrapping at 28pt."""
    if kicker:
        add_text(slide, Inches(0.6), Inches(0.4), Inches(10), Inches(0.35),
                 kicker.upper(), size=12, bold=True, color=PRIMARY)
    add_text(slide, Inches(0.6), Inches(0.72), Inches(12.1), Inches(1.2),
             title, size=28, bold=True, color=DARKER, line_spacing=1.1)


def chart_slide(title, kicker, chart_path, commentary, chart_height_in=4.8,
                chart_top_in=2.1):
    """Content slide with page header, a chart image, and a right-side commentary panel."""
    slide = prs.slides.add_slide(BLANK)
    add_rect(slide, 0, 0, SW, SH, LIGHT)
    add_page_header(slide, title, kicker)

    # Chart on the left, sized by height
    from PIL import Image
    img = Image.open(chart_path)
    ar = img.width / img.height
    ch_h = Inches(chart_height_in)
    ch_w = Emu(int(ch_h * ar))
    # Don't exceed 8" wide
    max_w = Inches(8.2)
    if ch_w > max_w:
        ch_w = max_w
        ch_h = Emu(int(ch_w / ar))
    ch_left = Inches(0.55)
    ch_top = Inches(chart_top_in)
    slide.shapes.add_picture(str(chart_path), ch_left, ch_top, width=ch_w, height=ch_h)

    # Commentary panel on the right
    panel_left = Inches(9.05)
    panel_top = Inches(2.1)
    panel_w = Inches(3.75)
    panel_h = Inches(5.4)
    panel = add_rect(slide, panel_left, panel_top, panel_w, panel_h, WHITE)
    panel.line.color.rgb = RGBColor(0xE1, 0xE7, 0xED)
    panel.line.width = Pt(0.75)

    # Takeaway kicker
    add_text(slide, panel_left + Inches(0.3), panel_top + Inches(0.3),
             panel_w - Inches(0.6), Inches(0.3),
             "TAKEAWAY", size=11, bold=True, color=ACCENT)

    # Commentary body
    add_bullets(slide, panel_left + Inches(0.3), panel_top + Inches(0.65),
                panel_w - Inches(0.6), panel_h - Inches(0.9),
                commentary, size=13, line_spacing=1.3)

    # Slide number in bottom right
    return slide


def add_slide_number(slide, n, total=10):
    add_text(slide, Inches(12.3), Inches(7.1), Inches(1.0), Inches(0.3),
             f"{n} / {total}", size=10, color=TEXT_MUTED, align=PP_ALIGN.RIGHT)


# ===========================================================================
# Slide 1 — Title (dark)
# ===========================================================================
s = prs.slides.add_slide(BLANK)
add_rect(s, 0, 0, SW, SH, DARK)
# Accent bar left
add_rect(s, 0, 0, Inches(0.25), SH, ACCENT)

add_text(s, Inches(0.9), Inches(1.2), Inches(12), Inches(0.5),
         "COMP 560 · FINAL PROJECT", size=14, bold=True, color=RGBColor(0xAE, 0xCC, 0xDF))

add_text(s, Inches(0.9), Inches(1.8), Inches(12), Inches(2.0),
         "Object Re-Identification", size=54, bold=True, color=WHITE,
         line_spacing=1.05)

add_text(s, Inches(0.9), Inches(3.6), Inches(12), Inches(1.2),
         "Cross-domain feature learning with DINOv2,\nBNNeck, and two-level PK sampling",
         size=24, color=RGBColor(0xCA, 0xDC, 0xFC), line_spacing=1.2)

# Stat strip at bottom
add_text(s, Inches(0.9), Inches(5.6), Inches(12), Inches(0.3),
         "KEY RESULTS", size=11, bold=True, color=ACCENT)

def stat(x, big, small, color=WHITE):
    add_text(s, Inches(x), Inches(5.95), Inches(3.0), Inches(0.9),
             big, size=40, bold=True, color=color)
    add_text(s, Inches(x), Inches(6.55), Inches(3.0), Inches(0.4),
             small, size=13, color=RGBColor(0xCA, 0xDC, 0xFC))

stat(0.9, "99.95%", "Rank-1 on Dataset A")
stat(4.0, "88.19%", "mAP on Dataset A")
stat(7.0, "+12.4", "Cross-domain mAP gain")
stat(10.1, "255", "images/sec on MPS")

# Footer
add_text(s, Inches(0.9), Inches(7.1), Inches(12), Inches(0.3),
         "Josh Robertson  ·  Student ID 730711465  ·  April 2026",
         size=11, color=RGBColor(0x8A, 0xA5, 0xBB))


# ===========================================================================
# Slide 2 — Task + Datasets
# ===========================================================================
s = prs.slides.add_slide(BLANK)
add_rect(s, 0, 0, SW, SH, LIGHT)
add_page_header(s, "The task: identify individuals across views", "PROBLEM")

# Two-column: left = formal task, right = two dataset cards
# Left column: task
lx = Inches(0.6); ly = Inches(2.15); lw = Inches(5.6)
add_text(s, lx, ly, lw, Inches(0.35),
         "Embedding learning for retrieval", size=16, bold=True, color=PRIMARY)

add_bullets(s, lx, ly + Inches(0.45), lw, Inches(4.5), [
    "Encode an image into a discriminative embedding.",
    "Rank a gallery by cosine similarity to a query.",
    "Score by Rank-K and mean average precision (mAP).",
    "Submitted interface: StudentModel.encode(images) → (B, 256) L2-normed.",
    "Evaluation is efficiency-aware (throughput, memory, embed dim).",
], size=15)

# Right: dataset cards
rx = Inches(6.6); rw = Inches(6.1)
# Card A
card_a = add_rect(s, rx, Inches(2.1), rw, Inches(2.35), WHITE)
card_a.line.color.rgb = PRIMARY
card_a.line.width = Pt(1.25)
add_rect(s, rx, Inches(2.1), Inches(0.18), Inches(2.35), PRIMARY)
add_text(s, rx + Inches(0.45), Inches(2.22), rw - Inches(0.7), Inches(0.4),
         "DATASET A  ·  OpenAnimals", size=15, bold=True, color=DARKER)
add_text(s, rx + Inches(0.45), Inches(2.62), rw - Inches(0.7), Inches(0.4),
         "Wildlife, 37 sub-datasets, available now", size=12, color=TEXT_MUTED)
add_bullets(s, rx + Inches(0.45), Inches(3.05), rw - Inches(0.7), Inches(1.35), [
    "109 927 train images · 30 561 test images",
    "10 249 identities · 37 species / sub-datasets",
    "Query/gallery split built from test.parquet",
], size=12, line_spacing=1.35)

# Card B
card_b = add_rect(s, rx, Inches(4.6), rw, Inches(2.35), WHITE)
card_b.line.color.rgb = ACCENT
card_b.line.width = Pt(1.25)
add_rect(s, rx, Inches(4.6), Inches(0.18), Inches(2.35), ACCENT)
add_text(s, rx + Inches(0.45), Inches(4.72), rw - Inches(0.7), Inches(0.4),
         "DATASET B  ·  Vehicle surveillance", size=15, bold=True, color=DARKER)
add_text(s, rx + Inches(0.45), Inches(5.12), rw - Inches(0.7), Inches(0.4),
         "Withheld — evaluated at grading time", size=12, color=ACCENT)
add_bullets(s, rx + Inches(0.45), Inches(5.55), rw - Inches(0.7), Inches(1.3), [
    "20 cameras · ~776 identities · ~13 k images",
    "Same-camera matches excluded at eval",
    "Different domain → model must generalize",
], size=12, line_spacing=1.35)

add_slide_number(s, 2)


# ===========================================================================
# Slide 3 — Approach overview
# ===========================================================================
s = prs.slides.add_slide(BLANK)
add_rect(s, 0, 0, SW, SH, LIGHT)
add_page_header(s, "Approach: one model, built for two domains", "METHOD")

# Pipeline: Input → Backbone → Head → Loss
# 4 boxes in a row with arrows
y = Inches(2.15); h = Inches(2.4)
box_w = Inches(2.75); gap = Inches(0.25)
start_x = Inches(0.6)

def pipeline_box(x, title, body, color=PRIMARY):
    # Colored header bar
    add_rect(s, x, y, box_w, Inches(0.55), color)
    add_text(s, x + Inches(0.15), y + Inches(0.10), box_w - Inches(0.3), Inches(0.4),
             title, size=13, bold=True, color=WHITE)
    # White body
    body_box = add_rect(s, x, y + Inches(0.55), box_w, h - Inches(0.55), WHITE)
    body_box.line.color.rgb = RGBColor(0xE1, 0xE7, 0xED)
    body_box.line.width = Pt(0.5)
    add_bullets(s, x + Inches(0.15), y + Inches(0.7), box_w - Inches(0.3), h - Inches(0.8),
                body, size=11, line_spacing=1.35)

pipeline_box(start_x + (box_w + gap) * 0, "1 · BACKBONE",
             ["DINOv2 ViT-S/14",
              "SSL pre-train on LVD-142M",
              "21 M params",
              "Freeze blocks 0–5 (6 of 12)"])
pipeline_box(start_x + (box_w + gap) * 1, "2 · HEAD",
             ["BNNeck (Luo et al. 2019)",
              "Linear 384 → 256",
              "BatchNorm1d, no bias",
              "Pre-BN → triplet,\npost-BN → classifier"])
pipeline_box(start_x + (box_w + gap) * 2, "3 · SAMPLER",
             ["Two-level PK",
              "Uniform sub-dataset",
              "P = 16 ids × K = 4 imgs",
              "Batch 64 · harder negatives"])
pipeline_box(start_x + (box_w + gap) * 3, "4 · LOSS", color=ACCENT,
             body=["CosFace (m = 0.35, s = 30)",
                   "+ hard-mining triplet",
                   "Label smoothing 0.1",
                   "AdamW · cosine LR · bf16"])

# Bottom strip: three validation / training callouts
y2 = Inches(5.05)
cx = Inches(0.6)
cw = Inches(4.0); cg = Inches(0.2)
def callout(x, title, body, color=PRIMARY):
    box = add_rect(s, x, y2, cw, Inches(1.9), WHITE)
    box.line.color.rgb = color
    box.line.width = Pt(1.0)
    add_rect(s, x, y2, cw, Inches(0.1), color)
    add_text(s, x + Inches(0.2), y2 + Inches(0.25), cw - Inches(0.4), Inches(0.35),
             title, size=13, bold=True, color=DARKER)
    add_text(s, x + Inches(0.2), y2 + Inches(0.7), cw - Inches(0.4), Inches(1.2),
             body, size=12, color=TEXT_DARK, line_spacing=1.3)

callout(cx + (cw + cg) * 0, "Partial freeze",
        "Preserve the DINOv2 features that transfer to Dataset B. "
        "Only the top half of the transformer moves.")
callout(cx + (cw + cg) * 1, "Two-level sampling",
        "Balanced across 37 sub-datasets → rare domains aren't drowned by "
        "CatIndividualImages.", color=ACCENT)
callout(cx + (cw + cg) * 2, "Cross-domain dev set",
        "3 held-out sub-datasets act as a proxy for Dataset B. Save on dev "
        "mAP, not train loss.")

add_slide_number(s, 3)


# ===========================================================================
# Slide 4 — Training losses (chart 01)
# ===========================================================================
s = chart_slide(
    "Training losses decreased cleanly over 25 epochs",
    "CONVERGENCE",
    CHARTS / "01_training_losses.png",
    [
        "CosFace dropped 16.95 → 3.08 over 25 epochs.",
        "Triplet loss fell from 0.35 → 0.08 — hardest-positive and hardest-negative pairs became well separated.",
        "No instability, no NaN restarts. bf16 autocast on MPS was stable throughout.",
        "~19 min per epoch on Apple Silicon MPS (total ≈ 8 h).",
    ],
    chart_height_in=4.8, chart_top_in=2.1,
)
add_slide_number(s, 4)


# ===========================================================================
# Slide 5 — Cross-domain dev mAP (chart 02)  [KEY SLIDE]
# ===========================================================================
s = chart_slide(
    "Cross-domain generalization improved — and we kept the best checkpoint",
    "KEY SIGNAL — DATASET B PROXY",
    CHARTS / "02_cross_domain_dev_map.png",
    [
        "Three sub-datasets held out from training, evaluated every 2 epochs.",
        "Mean mAP peaked at epoch 10 (51.05, +12.4 vs zero-shot).",
        "Later epochs regressed slightly (50.0 by epoch 25) — classifier starts memorizing.",
        "Checkpoint saver held epoch 10 automatically. This is the weights we ship.",
    ],
    chart_height_in=4.8, chart_top_in=2.1,
)
add_slide_number(s, 5)


# ===========================================================================
# Slide 6 — Dataset A results (chart 03)
# ===========================================================================
s = chart_slide(
    "Dataset A results — crushed the baseline",
    "PERFORMANCE",
    CHARTS / "03_dataset_a_vs_baseline.png",
    [
        "Rank-1: 99.95 %  (vs VLM baseline ~60 %)",
        "mAP: 88.19 %  (vs VLM baseline ~50 %)",
        "All 7 334 queries found a true match by Rank-5.",
        "Combined score 94.07 % · mINP 34.62 %.",
    ],
    chart_height_in=4.8, chart_top_in=2.1,
)
add_slide_number(s, 6)


# ===========================================================================
# Slide 7 — Efficiency (chart 04)
# ===========================================================================
s = chart_slide(
    "Efficiency — small, fast, cheap",
    "EFFICIENCY · 30% OF GRADE",
    CHARTS / "04_efficiency.png",
    [
        "255 img/s on MPS bf16 — 1 275× the VLM baseline.",
        "~1.5 GB peak memory — 40× less than the VLM.",
        "256-dim embeddings — 1.5× smaller than VLM's 384.",
        "Projected ≥ 700 img/s on a consumer CUDA GPU.",
    ],
    chart_height_in=4.6, chart_top_in=2.05,
)
add_slide_number(s, 7)


# ===========================================================================
# Slide 8 — Backbone benchmark (chart 05)
# ===========================================================================
s = chart_slide(
    "Backbone choice: DINOv2 over ConvNeXt — on purpose",
    "DAY-1 BENCHMARK",
    CHARTS / "05_backbone_benchmark.png",
    [
        "Both backbones cleared the 30 img/s efficiency floor easily.",
        "ConvNeXt-Tiny was ~12 % faster but uses ImageNet-supervised features.",
        "DINOv2's self-supervised pre-training transfers far better to unseen Dataset B.",
        "We traded a small throughput hit for a large generalization win.",
    ],
    chart_height_in=4.6, chart_top_in=2.05,
)
add_slide_number(s, 8)


# ===========================================================================
# Slide 9 — Dataset composition (chart 06, portrait)
# ===========================================================================
s = prs.slides.add_slide(BLANK)
add_rect(s, 0, 0, SW, SH, LIGHT)
add_page_header(s, "37 sub-datasets, wildly imbalanced — why sampling matters", "DATA")

# Put the portrait chart on the right, commentary on the left
img_path = CHARTS / "06_dataset_composition.png"
from PIL import Image
img = Image.open(img_path)
ar = img.width / img.height  # ~1.15
ch_h = Inches(5.0)
ch_w = Emu(int(ch_h * ar))
ch_left = Inches(8.3)
ch_top = Inches(2.1)
# Shift if too wide
if ch_left + ch_w > SW - Inches(0.3):
    ch_w = SW - Inches(0.3) - ch_left
    ch_h = Emu(int(ch_w / ar))
s.shapes.add_picture(str(img_path), ch_left, ch_top, width=ch_w, height=ch_h)

# Left side commentary + stat
lx = Inches(0.6); lw = Inches(7.4)
add_text(s, lx, Inches(2.15), lw, Inches(0.35),
         "Largest is 44× the smallest", size=16, bold=True, color=PRIMARY)

add_bullets(s, lx, Inches(2.65), lw, Inches(3.1), [
    "CatIndividualImages: 10 399 images — 44× the smallest sub-dataset.",
    "ReunionTurtles: 237 images. ZakynthosTurtles: 114.",
    "Naive PK sampling over 109 k images would draw a cat identity ~90 %  of batches and never see the rare sub-datasets.",
    "Our two-level sampler draws the sub-dataset uniformly first, then P=16 identities within it — rare domains get equal batches.",
    "Bonus: same-domain batches produce harder negatives for the triplet loss (cat vs cat, not cat vs whale).",
], size=13, line_spacing=1.4)

# Legend callout for the held-out subs
lg = add_rect(s, lx, Inches(5.8), lw, Inches(1.35), WHITE)
lg.line.color.rgb = ACCENT
lg.line.width = Pt(1.0)
add_text(s, lx + Inches(0.25), Inches(5.95), lw - Inches(0.5), Inches(0.4),
         "HELD OUT AS CROSS-DOMAIN DEV", size=11, bold=True, color=ACCENT)
add_text(s, lx + Inches(0.25), Inches(6.3), lw - Inches(0.5), Inches(0.9),
         "AAUZebraFish (fish)  ·  PolarBearVidID (video)  ·  SMALST (synthetic) — "
         "three visually distinct sub-datasets chosen as proxies for the unseen Dataset B.",
         size=13, color=TEXT_DARK, line_spacing=1.3)

add_slide_number(s, 9)


# ===========================================================================
# Slide 10 — Conclusion (dark)
# ===========================================================================
s = prs.slides.add_slide(BLANK)
add_rect(s, 0, 0, SW, SH, DARK)
add_rect(s, 0, 0, Inches(0.25), SH, ACCENT)

add_text(s, Inches(0.9), Inches(0.6), Inches(12), Inches(0.4),
         "TAKEAWAYS", size=14, bold=True, color=RGBColor(0xAE, 0xCC, 0xDF))

add_text(s, Inches(0.9), Inches(1.1), Inches(12), Inches(1.0),
         "Cross-domain ReID, done right", size=40, bold=True, color=WHITE)

# 3 columns of stats
col_y = Inches(2.5); col_h = Inches(3.2); col_w = Inches(3.9); col_g = Inches(0.2)
col_x0 = Inches(0.9)
def concl_col(x, kicker, headline, body):
    box = add_rect(s, x, col_y, col_w, col_h, DARKER)
    box.line.color.rgb = RGBColor(0x23, 0x4A, 0x62)
    box.line.width = Pt(0.75)
    add_rect(s, x, col_y, col_w, Inches(0.08), ACCENT)
    add_text(s, x + Inches(0.3), col_y + Inches(0.35), col_w - Inches(0.6), Inches(0.3),
             kicker, size=11, bold=True, color=ACCENT)
    add_text(s, x + Inches(0.3), col_y + Inches(0.7), col_w - Inches(0.6), Inches(0.7),
             headline, size=22, bold=True, color=WHITE, line_spacing=1.15)
    add_text(s, x + Inches(0.3), col_y + Inches(1.9), col_w - Inches(0.6), Inches(1.2),
             body, size=13, color=RGBColor(0xCA, 0xDC, 0xFC), line_spacing=1.4)

concl_col(col_x0 + (col_w + col_g) * 0, "PERFORMANCE",
          "99.95 / 88.19",
          "Rank-1 / mAP on Dataset A.\n+40 / +38 over the provided\nVLM baseline.")

concl_col(col_x0 + (col_w + col_g) * 1, "GENERALIZATION",
          "+12.4 mean mAP",
          "On 3 held-out sub-datasets.\nPartial freeze + SSL backbone\nkept Dataset B transfer intact.")

concl_col(col_x0 + (col_w + col_g) * 2, "EFFICIENCY",
          "1 275× faster",
          "255 img/s at 1.5 GB mem,\n256-dim embeddings. Well\nunder every efficiency cap.")

# Footer bar
add_text(s, Inches(0.9), Inches(6.35), Inches(12), Inches(0.4),
         "DINOv2 ViT-S/14  ·  BNNeck  ·  CosFace + TripletHard  ·  two-level PK sampler",
         size=14, bold=True, color=RGBColor(0xAE, 0xCC, 0xDF), align=PP_ALIGN.CENTER)

add_text(s, Inches(0.9), Inches(6.85), Inches(12), Inches(0.3),
         "Trained on Apple Silicon MPS · 25 epochs · best checkpoint at epoch 10",
         size=12, color=RGBColor(0x8A, 0xA5, 0xBB), align=PP_ALIGN.CENTER)

# No slide number on the conclusion for visual cleanliness

# ---- Save ----
prs.save(str(OUT))
print(f"Saved {OUT}  ({OUT.stat().st_size / 1024:.0f} KB, {len(prs.slides)} slides)")
